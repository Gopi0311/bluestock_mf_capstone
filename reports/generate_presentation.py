"""
generate_presentation.py - Day 7: Create 12-slide presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUTPUT = Path(__file__).parent / "Bluestock_MF_Presentation.pptx"
REPORTS = Path(__file__).parent

BLUE = RGBColor(0x1B, 0x4F, 0x9B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF5, 0xA6, 0x23)

def add_slide(prs, layout=1):
    return prs.slides.add_slide(prs.slide_layouts[layout])

def set_title(slide, text, color=None):
    title = slide.shapes.title
    title.text = text
    if color:
        title.text_frame.paragraphs[0].runs[0].font.color.rgb = color
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    title.text_frame.paragraphs[0].runs[0].font.bold = True

def add_text(slide, text, left, top, width, height, size=16, bold=False, color=None):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # SLIDE 1 - Title
    slide = add_slide(prs, 0)
    slide.shapes.title.text = "Bluestock Fintech"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = BLUE
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(40)
    slide.placeholders[1].text = "Mutual Fund Analytics Platform\nCapstone Project | June 2026"
    add_text(slide, "Upputholla Gopi | Data Analyst Intern | VVIT Guntur",
             1, 6, 10, 0.5, size=14, color=BLUE)

    # SLIDE 2 - Problem & Objective
    slide = add_slide(prs, 1)
    set_title(slide, "Problem & Objective", BLUE)
    content = """PROBLEMS:
• Data Fragmentation — NAV, AUM, SIP data in different formats
• No unified platform for risk-adjusted fund comparison
• Retail investors cannot track benchmark performance
• Static PDF reports — no real-time insights

OBJECTIVES:
• Build automated ETL pipeline from AMFI public data
• Design normalized SQL database (star schema)
• Compute Sharpe, Alpha, Beta, VaR for 40 funds
• Build interactive Power BI dashboard"""
    slide.placeholders[1].text = content

    # SLIDE 3 - Data Sources
    slide = add_slide(prs, 1)
    set_title(slide, "Data Sources & Scale", BLUE)
    slide.placeholders[1].text = """DATA SOURCES:
• AMFI India — NAV, AUM, SIP, Folio data
• mfapi.in — Live NAV API (no auth required)
• NSE/BSE — Benchmark index prices

SCALE:
• 40 Mutual Fund Schemes across 10 AMCs
• 46,000+ Daily NAV records (2022-2026)
• 32,000+ Investor Transactions
• 8,050 Benchmark Index records
• Total: 87,000+ rows of financial data"""

    # SLIDE 4 - Architecture
    slide = add_slide(prs, 1)
    set_title(slide, "System Architecture", BLUE)
    slide.placeholders[1].text = """ETL PIPELINE:

EXTRACT → AMFI APIs, mfapi.in, CSV datasets
      ↓
TRANSFORM → Python/Pandas: clean, validate, compute metrics
      ↓
LOAD → SQLite database (5-table star schema)
      ↓
ANALYSE → Jupyter notebooks: EDA, performance metrics
      ↓
VISUALISE → Power BI: 4-page interactive dashboard"""

    # SLIDE 5 - EDA Highlights 1
    slide = add_slide(prs, 1)
    set_title(slide, "EDA Highlights — Market Trends", BLUE)
    slide.placeholders[1].text = """KEY FINDINGS:

📈 NAV Trend: All 40 funds grew consistently 2022-2026
   Strong 2023 bull run, minor 2024 corrections

🏦 SBI Dominance: Rs.12.5 Lakh Crore AUM
   Nearly 2x the second-largest AMC (ICICI Prudential)

💰 SIP All-time High: Rs.31,002 Crore in Dec 2025
   Reflects India's growing retail investment culture

📊 Folio Growth: 13.26 Cr → 26.12 Cr (97% growth in 4 years)"""

    # SLIDE 6 - EDA Highlights 2
    slide = add_slide(prs, 1)
    set_title(slide, "EDA Highlights — Investor Behaviour", BLUE)
    slide.placeholders[1].text = """INVESTOR INSIGHTS:

👥 Demographics: 26-35 age group = largest investor segment
   Young professionals driving MF growth in India

🗺️ Geography: Maharashtra & Karnataka lead investments
   T30 cities = 70% of total AUM

⚖️ Gender Gap: Male 65% vs Female 35% investors
   Opportunity for women-focused MF campaigns

🏙️ B30 Growth: Beyond Top-30 cities growing faster in SIPs"""

    # SLIDE 7 - Performance Metrics
    slide = add_slide(prs, 1)
    set_title(slide, "Performance Metrics — Methodology", BLUE)
    slide.placeholders[1].text = """METRICS COMPUTED FOR ALL 40 FUNDS:

• CAGR (1yr, 3yr, 5yr) — Annualised compound returns
• Sharpe Ratio — Risk-adjusted return vs RBI repo rate (6.5%)
• Sortino Ratio — Penalises only downside volatility
• Alpha & Beta — OLS regression vs Nifty 100 benchmark
• Max Drawdown — Worst peak-to-trough decline
• VaR (95%) — Daily loss threshold at 95% confidence
• CVaR — Expected loss beyond VaR threshold
• Tracking Error — Std dev of fund vs benchmark returns"""

    # SLIDE 8 - Performance Results
    slide = add_slide(prs, 1)
    set_title(slide, "Performance Analytics — Results", BLUE)
    slide.placeholders[1].text = """FUND SCORECARD (Composite 0-100 Score):
Weighted: 30% Return + 25% Sharpe + 20% Alpha + 15% Expense + 10% Max DD

TOP FINDINGS:
• Funds with Sharpe > 1 consistently outperform benchmark
• Large cap funds show Beta close to 1.0 (market-tracking)
• Small/Mid cap funds have higher Alpha but higher VaR
• Liquid funds show lowest VaR — best for capital preservation
• High HHI funds concentrated in Consumer Goods & IT sectors

RISK INSIGHTS:
• Rolling Sharpe declined in 2024 — market-wide risk increase
• At-risk SIP investors (gap > 35 days) need re-engagement"""

    # SLIDE 9 - Dashboard Page 1&2
    slide = add_slide(prs, 1)
    set_title(slide, "Power BI Dashboard — Pages 1 & 2", BLUE)
    slide.placeholders[1].text = """PAGE 1 — INDUSTRY OVERVIEW:
• KPI Cards: Total AUM, SIP Inflow, Schemes, Folios
• Bar Chart: AUM by Fund House (SBI dominance visible)
• Line Chart: Monthly SIP Inflow Trend 2022-2025
• Slicer: Fund Category filter

PAGE 2 — FUND PERFORMANCE:
• Scatter Plot: Return vs Risk (daily_return vs NAV)
• Fund Scorecard Table: sortable by all metrics
• Fund House Slicer: filter by AMC"""

    # SLIDE 10 - Dashboard Page 3&4
    slide = add_slide(prs, 1)
    set_title(slide, "Power BI Dashboard — Pages 3 & 4", BLUE)
    slide.placeholders[1].text = """PAGE 3 — INVESTOR ANALYTICS:
• Bar Chart: Transaction Amount by State
• Donut Chart: SIP vs Lumpsum vs Redemption split
• Bar Chart: Age Group vs Investment Amount
• Monthly Transaction Volume Line Chart
• Slicers: State, City Tier (T30/B30)

PAGE 4 — SIP & MARKET TRENDS:
• SIP Inflow Monthly Trend
• Category Inflows Bar Chart
• KPI: Active SIP Accounts, SIP AUM
• Month Slicer for time filtering"""

    # SLIDE 11 - Key Findings
    slide = add_slide(prs, 1)
    set_title(slide, "Key Findings & Recommendations", BLUE)
    slide.placeholders[1].text = """KEY FINDINGS:
1. SBI leads industry with Rs.12.5L Cr AUM
2. SIP culture growing — Rs.31,002 Cr monthly high
3. Young investors (26-35) driving growth
4. Large cap funds best risk-adjusted returns
5. B30 cities — untapped growth opportunity

RECOMMENDATIONS:
• Low risk investors → Liquid/Short-term Debt funds
• Moderate risk → Large Cap with Sharpe > 1
• High risk → Mid/Small Cap with monitored drawdown
• AMCs should target B30 cities for SIP campaigns
• Re-engage at-risk SIP investors (gap > 35 days)"""

    # SLIDE 12 - Thank You
    slide = add_slide(prs, 0)
    slide.shapes.title.text = "Thank You!"
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = BLUE
    slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(44)
    slide.placeholders[1].text = "Upputholla Gopi\nData Analyst Intern — Bluestock Fintech\nVVIT, Guntur | June 2026\nGitHub: github.com/Gopi0311/bluestock_mf_capstone"

    prs.save(str(OUTPUT))
    print(f"✅ Presentation saved: {OUTPUT}")

if __name__ == "__main__":
    build_presentation()